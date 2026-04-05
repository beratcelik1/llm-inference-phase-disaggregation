"""Clean minimalist PPTX presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Minimal palette
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
BORDER_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
BLUE = RGBColor(0x1A, 0x56, 0xDB)  # primary accent
NAVY = RGBColor(0x0F, 0x17, 0x2A)  # dark slides
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)  # warnings/problems
GREEN = RGBColor(0x1E, 0x8A, 0x4C)  # positive results
TEAL = RGBColor(0x0E, 0x6E, 0x6E)  # secondary


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE

    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def box(slide, x, y, w, h):
    return slide.shapes.add_textbox(x, y, w, h)


def txt(tf, text, sz=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Palatino Linotype"
    p.alignment = align
    return p


def para(tf, text, sz=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, sp=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Palatino Linotype"
    p.alignment = align
    p.space_before = sp
    return p


def bullet(tf, text, sz=16, color=DARK, level=0, bold=False, sp=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Palatino Linotype"
    p.level = level
    p.space_before = sp
    return p


def tbl(slide, rows, cols, x, y, w, h):
    return slide.shapes.add_table(rows, cols, x, y, w, h).table


def cell(table, r, c, text, sz=14, bold=False, color=DARK, fill=None):
    cl = table.cell(r, c)
    cl.text = ""
    p = cl.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Palatino Linotype"
    if fill:
        cl.fill.solid()
        cl.fill.fore_color.rgb = fill


L = Inches(1.2)  # left margin
T_TITLE = Inches(0.5)
T_BODY = Inches(1.6)
W = Inches(11)

FIGURES = Path(__file__).parent / "figures"


def fig_slide(title: str, img_name: str, caption: str = "") -> None:
    """Add a slide showing a paper figure."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    tb = box(s, L, Inches(0.3), W, Inches(0.6))
    txt(tb.text_frame, title, sz=24, bold=True, color=BLACK)
    img_path = FIGURES / img_name
    if img_path.exists():
        pic = s.shapes.add_picture(
            str(img_path), Inches(1.9), Inches(1.2), width=Inches(9.5)
        )
        if caption:
            cy = pic.top + pic.height + Inches(0.15)
            tb = box(s, Inches(1.9), cy, Inches(9.5), Inches(0.4))
            txt(tb.text_frame, caption, sz=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ====== SLIDE 1: TITLE ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
rect(s, Inches(0), Inches(3.4), Inches(13.333), Pt(2), BLUE)

tb = box(s, L, Inches(1.8), W, Inches(1.2))
txt(
    tb.text_frame,
    "Phase Disaggregation for LLM Inference",
    sz=42,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.LEFT,
)

tb = box(s, L, Inches(3.7), W, Inches(0.6))
txt(
    tb.text_frame,
    "Splitwise  &  DistServe",
    sz=26,
    color=RGBColor(0x88, 0xAA, 0xDD),
    align=PP_ALIGN.LEFT,
)

tb = box(s, L, Inches(5.0), W, Inches(1.2))
txt(
    tb.text_frame,
    "ECE 5545: Machine Learning Hardware & Systems",
    sz=18,
    color=LIGHT_GRAY,
)
para(
    tb.text_frame,
    "Berat Celik  &  Jiayang (Ethan) Chen",
    sz=20,
    bold=True,
    color=WHITE,
    sp=Pt(12),
)
para(tb.text_frame, "Spring 2026", sz=16, color=LIGHT_GRAY, sp=Pt(8))

# ====== SLIDE 2: OUTLINE ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "Outline", sz=34, bold=True, color=BLACK)

# Part 1
tb = box(s, L, Inches(1.8), Inches(5.2), Inches(4))
txt(tb.text_frame, "Part 1: Berat", sz=22, bold=True, color=BLUE)
bullet(tb.text_frame, "LLM inference background", sz=17, color=DARK, sp=Pt(14))
bullet(tb.text_frame, "The two phases and why they differ", sz=17, color=DARK)
bullet(tb.text_frame, "The colocation problem", sz=17, color=DARK)
bullet(tb.text_frame, "Splitwise: production insights", sz=17, color=DARK)
bullet(tb.text_frame, "Architecture and KV-cache transfer", sz=17, color=DARK)
bullet(tb.text_frame, "Heterogeneous cluster designs", sz=17, color=DARK)
bullet(tb.text_frame, "Evaluation", sz=17, color=DARK)

# Part 2
tb = box(s, Inches(7), Inches(1.8), Inches(5.2), Inches(4))
txt(tb.text_frame, "Part 2: Ethan", sz=22, bold=True, color=TEAL)
bullet(tb.text_frame, "DistServe: goodput optimization", sz=17, color=DARK, sp=Pt(14))
bullet(tb.text_frame, "Formal tradeoff analysis", sz=17, color=DARK)
bullet(tb.text_frame, "Placement algorithms", sz=17, color=DARK)
bullet(tb.text_frame, "Online scheduling", sz=17, color=DARK)
bullet(tb.text_frame, "Evaluation", sz=17, color=DARK)
bullet(tb.text_frame, "Comparison and future directions", sz=17, color=DARK)

# ====== SLIDE 3: HOW LLM INFERENCE WORKS ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "How LLM Inference Works", sz=34, bold=True, color=BLACK)

# Phase 1
rect(s, L, Inches(1.7), Inches(5.2), Inches(2.8), OFF_WHITE)
tb = box(s, Inches(1.5), Inches(1.8), Inches(4.6), Inches(2.6))
txt(tb.text_frame, "Phase 1: Prefill", sz=22, bold=True, color=BLUE)
bullet(
    tb.text_frame,
    "Processes all input tokens in parallel",
    sz=16,
    color=DARK,
    sp=Pt(10),
)
bullet(tb.text_frame, "One forward pass through the model", sz=16, color=DARK)
bullet(tb.text_frame, "Produces the first output token", sz=16, color=DARK)
bullet(tb.text_frame, "Generates the KV-cache (stored context)", sz=16, color=DARK)
bullet(tb.text_frame, "Compute-bound", sz=16, color=BLUE, bold=True)

# Phase 2
rect(s, Inches(6.9), Inches(1.7), Inches(5.2), Inches(2.8), OFF_WHITE)
tb = box(s, Inches(7.2), Inches(1.8), Inches(4.6), Inches(2.6))
txt(tb.text_frame, "Phase 2: Decoding", sz=22, bold=True, color=TEAL)
bullet(tb.text_frame, "Generates tokens one at a time", sz=16, color=DARK, sp=Pt(10))
bullet(tb.text_frame, "Each step reads the full KV-cache", sz=16, color=DARK)
bullet(tb.text_frame, "Repeats until end of sequence", sz=16, color=DARK)
bullet(tb.text_frame, "Sequential, autoregressive", sz=16, color=DARK)
bullet(tb.text_frame, "Memory-bandwidth-bound", sz=16, color=TEAL, bold=True)

# Example
rect(s, L, Inches(4.9), W, Inches(1.0), OFF_WHITE)
tb = box(s, Inches(1.5), Inches(5.0), Inches(10.2), Inches(0.8))
txt(tb.text_frame, '"Is tomato a fruit?"', sz=17, bold=True, color=DARK)
para(
    tb.text_frame,
    'Prefill: process 5 tokens at once, output "Yes"    Decode: generate  ","  "it"  "is"  "."  one by one',
    sz=15,
    color=GRAY,
    sp=Pt(4),
)

tb = box(s, L, Inches(6.2), W, Inches(0.8))
txt(
    tb.text_frame,
    "KV-Cache: key/value tensors stored during prefill, read at every decode step. This is what gets transferred when phases are on different machines.",
    sz=14,
    color=LIGHT_GRAY,
)

# ====== SLIDE 4: CHARACTERISTICS TABLE ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "Two Phases, Very Different Profiles", sz=34, bold=True, color=BLACK)

table = tbl(s, 8, 3, Inches(1.5), Inches(1.6), Inches(10.3), Inches(4.5))
headers = ["", "Prefill", "Decoding"]
rows = [
    ["Bottleneck", "Compute-bound", "Memory bandwidth-bound"],
    ["Tokens per step", "All input tokens (parallel)", "One token (sequential)"],
    ["GPU utilization", "High", "Low without batching"],
    ["Power draw", "Near TDP", "Well below TDP"],
    ["Latency metric", "TTFT (time to first token)", "TPOT / TBT (time per token)"],
    ["Batching", "Limited benefit (already saturated)", "Significant benefit"],
    ["Preferred parallelism", "Tensor (intra-op)", "Pipeline (inter-op)"],
]
for i, h in enumerate(headers):
    cell(
        table,
        0,
        i,
        h,
        sz=15,
        bold=True,
        color=WHITE,
        fill=BLUE if i == 1 else (TEAL if i == 2 else RGBColor(0x55, 0x55, 0x55)),
    )
for r, row in enumerate(rows):
    for c, v in enumerate(row):
        f = OFF_WHITE if r % 2 == 0 else WHITE
        cell(table, r + 1, c, v, sz=13, color=DARK, fill=f, bold=(c == 0))

tb = box(s, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.5))
txt(
    tb.text_frame,
    "This asymmetry is what motivates both papers.",
    sz=18,
    bold=True,
    color=BLUE,
    align=PP_ALIGN.CENTER,
)

# ====== SLIDE 5: METRICS ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "Key Metrics", sz=34, bold=True, color=BLACK)

table = tbl(s, 7, 3, Inches(1.5), Inches(1.6), Inches(10.3), Inches(3.8))
h = ["Metric", "Definition", "Why it matters"]
d = [
    ["TTFT", "Time to first token", "User-perceived responsiveness"],
    ["TBT / TPOT", "Time between output tokens", "Streaming reading speed"],
    ["E2E Latency", "TTFT + TPOT x output length", "Total request time"],
    ["Throughput", "Requests per second", "System capacity"],
    ["Goodput", "Throughput while meeting SLO targets", "Cost efficiency"],
    ["SLO Attainment", "% of requests within latency bounds", "Service quality"],
]
for i, v in enumerate(h):
    cell(table, 0, i, v, sz=14, bold=True, color=WHITE, fill=RGBColor(0x55, 0x55, 0x55))
for r, row in enumerate(d):
    f = OFF_WHITE if r % 2 == 0 else WHITE
    for c, v in enumerate(row):
        cell(table, r + 1, c, v, sz=13, color=DARK, fill=f, bold=(c == 0))

rect(s, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.7), OFF_WHITE)
tb = box(s, Inches(1.8), Inches(5.75), Inches(9.7), Inches(0.6))
txt(
    tb.text_frame,
    'Example SLO: "90% of requests must have TTFT < 0.25s and TPOT < 0.1s"',
    sz=16,
    color=GRAY,
    align=PP_ALIGN.CENTER,
)

# ====== SLIDE 6: THE PROBLEM ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "The Problem with Colocation", sz=34, bold=True, color=RED_ACCENT)

# Three problems side by side
for i, (title, items) in enumerate(
    [
        (
            "Prefill/Decode Interference",
            [
                "One prefill in a decode batch:",
                "  3x latency hit (input=128)",
                "  5x latency hit (input=1024)",
                "Goes both directions",
            ],
        ),
        (
            "Parallelism Coupling",
            [
                "Prefill wants tensor parallelism",
                "Decoding wants pipeline parallelism",
                "Colocated: must pick one",
                "One phase always loses",
            ],
        ),
        (
            "Over-Provisioning",
            [
                "Must provision for worst case",
                "Colocated: 1.6 req/s (OPT-13B)",
                "Disaggregated: 3.3 req/s",
                "2.1x wasted capacity",
            ],
        ),
    ]
):
    x = Inches(0.8 + i * 4.15)
    rect(s, x, Inches(1.7), Inches(3.8), Inches(3.5), OFF_WHITE)
    tb = box(s, x + Inches(0.3), Inches(1.8), Inches(3.2), Inches(3.3))
    txt(tb.text_frame, title, sz=17, bold=True, color=RED_ACCENT)
    for item in items:
        bullet(tb.text_frame, item, sz=14, color=DARK, sp=Pt(6))

rect(s, L, Inches(5.6), W, Inches(0.8), BLUE)
tb = box(s, Inches(1.5), Inches(5.65), Inches(10.3), Inches(0.7))
txt(
    tb.text_frame,
    "Solution: put prefill and decoding on separate hardware.",
    sz=20,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)

fig_slide(
    "Prefill-Decode Interference",
    "ds_fig2_interference.png",
    "Source: Zhong et al., OSDI 2024, Figure 2",
)

# ====== SLIDE 7: SPLITWISE DIVIDER ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
rect(s, Inches(0), Inches(3.6), Inches(13.333), Pt(2), BLUE)

tb = box(s, L, Inches(2.2), W, Inches(1))
txt(tb.text_frame, "Splitwise", sz=48, bold=True, color=WHITE)

tb = box(s, L, Inches(4.0), W, Inches(1.2))
txt(
    tb.text_frame,
    "Efficient Generative LLM Inference Using Phase Splitting",
    sz=22,
    color=RGBColor(0x88, 0xAA, 0xDD),
)
para(
    tb.text_frame,
    "Patel et al. | University of Washington + Microsoft | ISCA 2024",
    sz=16,
    color=LIGHT_GRAY,
    sp=Pt(10),
)

# ====== SLIDE 8: PRODUCTION TRACES ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(
    tb.text_frame,
    "Insights from Azure Production Traces",
    sz=34,
    bold=True,
    color=BLACK,
)

tb = box(s, L, Inches(1.4), W, Inches(0.5))
txt(
    tb.text_frame,
    "Real traces from two Microsoft Azure LLM services (coding + conversation, Nov 2023)",
    sz=16,
    color=GRAY,
)

for i, (num, title, detail) in enumerate(
    [
        (
            "I",
            "Workloads vary widely",
            "Coding: 1500 prompt tokens, 13 output\nConversation: 1020 prompt, 129 output",
        ),
        (
            "II",
            "Token generation underutilizes GPUs",
            "60-70% of time spent with 20 or fewer\nactive tokens in the batch",
        ),
        (
            "III",
            "Most wall-clock time is in decoding",
            "Even large prompts (1500 tokens) take\nroughly the same time as 6 output tokens",
        ),
        (
            "IV",
            "Prompt batching saturates quickly",
            "Throughput plateaus after ~2048 tokens\nToken batching keeps scaling",
        ),
    ]
):
    x = Inches(0.7 + i * 3.15)
    rect(s, x, Inches(2.2), Inches(2.95), Inches(3.5), OFF_WHITE)
    tb = box(s, x + Inches(0.2), Inches(2.3), Inches(2.55), Inches(3.3))
    txt(tb.text_frame, f"Insight {num}", sz=14, bold=True, color=BLUE)
    para(tb.text_frame, title, sz=15, bold=True, color=BLACK, sp=Pt(6))
    for line in detail.split("\n"):
        para(tb.text_frame, line, sz=13, color=GRAY, sp=Pt(4))

tb = box(s, L, Inches(6.0), W, Inches(0.8))
txt(tb.text_frame, "Splitwise Figure 3-4, Table IV", sz=13, color=LIGHT_GRAY)

fig_slide(
    "Azure Production Trace Distributions",
    "sw_fig3_distributions.png",
    "Source: Patel et al., ISCA 2024, Figure 3",
)

# ====== SLIDE 9: HARDWARE INSIGHT ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "The Hardware Insight", sz=34, bold=True, color=BLACK)

table = tbl(s, 5, 4, Inches(1.5), Inches(1.5), Inches(6.5), Inches(2.5))
for i, v in enumerate(["", "A100", "H100", "Ratio"]):
    cell(table, 0, i, v, sz=14, bold=True, color=WHITE, fill=RGBColor(0x55, 0x55, 0x55))
for r, row in enumerate(
    [
        ["TTFT (coding)", "185 ms", "95 ms", "0.51x"],
        ["TBT (coding)", "52 ms", "31 ms", "0.70x"],
        ["Cost per request", "$0.42", "$0.52", "1.24x"],
        ["Energy", "1.37 Whr", "1.37 Whr", "1x"],
    ]
):
    f = OFF_WHITE if r % 2 == 0 else WHITE
    for c, v in enumerate(row):
        clr = (
            GREEN
            if (r == 0 and c == 3)
            else (RED_ACCENT if (r == 2 and c == 3) else DARK)
        )
        cell(table, r + 1, c, v, sz=13, color=clr, fill=f, bold=(c == 0))

rect(s, Inches(1.5), Inches(4.3), Inches(10.3), Inches(1.8), OFF_WHITE)
tb = box(s, Inches(1.8), Inches(4.4), Inches(9.7), Inches(1.6))
txt(tb.text_frame, "H100 has 3.43x more compute than A100", sz=17, color=DARK)
para(
    tb.text_frame,
    "Prefill benefits (TTFT: 0.51x). Token generation barely does (TBT: 0.70x).",
    sz=17,
    color=DARK,
    sp=Pt(8),
)
para(
    tb.text_frame,
    "Token generation can run on cheaper hardware without much performance loss.",
    sz=17,
    bold=True,
    color=BLUE,
    sp=Pt(8),
)

tb = box(s, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.5))
txt(
    tb.text_frame,
    "Power-capping token machines to 70% has virtually no latency impact.",
    sz=15,
    color=GRAY,
)

# ====== SLIDE 10: ARCHITECTURE ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "Splitwise Architecture", sz=34, bold=True, color=BLACK)

# Scheduler
rect(s, Inches(4.5), Inches(1.5), Inches(4.3), Inches(0.7), RGBColor(0x55, 0x55, 0x55))
tb = box(s, Inches(4.7), Inches(1.52), Inches(3.9), Inches(0.6))
txt(
    tb.text_frame,
    "Cluster-Level Scheduler (JSQ routing)",
    sz=15,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)

# Prompt Pool
rect(s, Inches(0.8), Inches(2.8), Inches(3.8), Inches(2.5), RGBColor(0xE8, 0xEE, 0xF8))
tb = box(s, Inches(1.0), Inches(2.9), Inches(3.4), Inches(2.3))
txt(tb.text_frame, "Prompt Pool", sz=20, bold=True, color=BLUE)
bullet(tb.text_frame, "Compute-optimized GPUs (H100)", sz=14, color=DARK, sp=Pt(10))
bullet(tb.text_frame, "Dedicated to prefill", sz=14, color=DARK)
bullet(tb.text_frame, "FCFS, batch up to 2048 tokens", sz=14, color=GRAY)

# Arrow
rect(s, Inches(4.8), Inches(3.5), Inches(3.7), Inches(0.5), OFF_WHITE)
tb = box(s, Inches(4.9), Inches(3.5), Inches(3.5), Inches(0.5))
txt(
    tb.text_frame,
    "KV-cache transfer  >>",
    sz=14,
    bold=True,
    color=GRAY,
    align=PP_ALIGN.CENTER,
)

# Token Pool
rect(s, Inches(8.7), Inches(2.8), Inches(3.8), Inches(2.5), RGBColor(0xE8, 0xF5, 0xF0))
tb = box(s, Inches(8.9), Inches(2.9), Inches(3.4), Inches(2.3))
txt(tb.text_frame, "Token Pool", sz=20, bold=True, color=TEAL)
bullet(tb.text_frame, "Cost-efficient GPUs (A100)", sz=14, color=DARK, sp=Pt(10))
bullet(tb.text_frame, "Dedicated to generation", sz=14, color=DARK)
bullet(tb.text_frame, "Continuous batching", sz=14, color=GRAY)

# Mixed Pool
rect(s, Inches(3.5), Inches(5.7), Inches(6.3), Inches(1.0), OFF_WHITE)
tb = box(s, Inches(3.7), Inches(5.75), Inches(5.9), Inches(0.9))
txt(tb.text_frame, "Mixed Pool", sz=17, bold=True, color=DARK)
para(
    tb.text_frame,
    "Flexible machines that handle overflow from either pool",
    sz=14,
    color=GRAY,
    sp=Pt(4),
)

fig_slide(
    "Splitwise System Architecture",
    "sw_fig10_architecture.png",
    "Source: Patel et al., ISCA 2024, Figure 10",
)

# ====== SLIDE 11: KV-CACHE TRANSFER ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "KV-Cache Transfer", sz=34, bold=True, color=BLACK)

# Naive
rect(s, L, Inches(1.7), Inches(5.2), Inches(2.8), RGBColor(0xFD, 0xF0, 0xEF))
tb = box(s, Inches(1.5), Inches(1.8), Inches(4.6), Inches(2.6))
txt(tb.text_frame, "Serialized (naive)", sz=20, bold=True, color=RED_ACCENT)
bullet(tb.text_frame, "Wait for full prefill to finish", sz=15, color=DARK, sp=Pt(10))
bullet(tb.text_frame, "Transfer entire KV-cache at once", sz=15, color=DARK)
bullet(
    tb.text_frame, "64% overhead on second token", sz=15, color=RED_ACCENT, bold=True
)

# Optimized
rect(s, Inches(6.9), Inches(1.7), Inches(5.2), Inches(2.8), RGBColor(0xEF, 0xFA, 0xF4))
tb = box(s, Inches(7.2), Inches(1.8), Inches(4.6), Inches(2.6))
txt(tb.text_frame, "Per-layer (optimized)", sz=20, bold=True, color=GREEN)
bullet(
    tb.text_frame,
    "Send each layer's cache as it computes",
    sz=15,
    color=DARK,
    sp=Pt(10),
)
bullet(tb.text_frame, "Transfer overlaps with computation", sz=15, color=DARK)
bullet(tb.text_frame, "16% overhead on second token", sz=15, color=GREEN, bold=True)

# Results
rect(s, L, Inches(4.9), W, Inches(1.8), OFF_WHITE)
tb = box(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.6))
txt(
    tb.text_frame,
    "Transfer overhead: 0.8% of end-to-end latency",
    sz=20,
    bold=True,
    color=BLUE,
)
para(
    tb.text_frame,
    "Constant 5-8ms regardless of prompt size (MSCCL++ zero-copy over InfiniBand)",
    sz=16,
    color=DARK,
    sp=Pt(8),
)
para(
    tb.text_frame,
    "Both papers confirm: KV-cache transfer is not the bottleneck.",
    sz=16,
    bold=True,
    color=GREEN,
    sp=Pt(8),
)

# ====== SLIDE 12: CLUSTER DESIGNS ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "Cluster Configurations", sz=34, bold=True, color=BLACK)

table = tbl(s, 5, 4, Inches(1.5), Inches(1.6), Inches(10.3), Inches(2.5))
for i, v in enumerate(["Configuration", "Prompt GPUs", "Token GPUs", "Tradeoff"]):
    cell(table, 0, i, v, sz=14, bold=True, color=WHITE, fill=RGBColor(0x55, 0x55, 0x55))
for r, row in enumerate(
    [
        ["Splitwise-AA", "DGX-A100", "DGX-A100", "Lowest cost"],
        ["Splitwise-HH", "DGX-H100", "DGX-H100", "Best performance"],
        ["Splitwise-HA", "DGX-H100", "DGX-A100", "Best cost-performance"],
        ["Splitwise-HHcap", "DGX-H100", "H100 at 70% power", "Power savings"],
    ]
):
    f = OFF_WHITE if r % 2 == 0 else WHITE
    for c, v in enumerate(row):
        cell(table, r + 1, c, v, sz=13, color=DARK, fill=f, bold=(c == 0))

# ====== SLIDE 13: SPLITWISE RESULTS ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "Splitwise Results", sz=34, bold=True, color=BLACK)

results = [
    ("Same power & cost budget", "2.15x more throughput vs Baseline-A100"),
    ("Same cost", "1.4x more throughput at 20% lower cost vs Baseline-H100"),
    ("Same throughput target", "25% lower power (HHcap) or 25% lower cost (AA)"),
]
for i, (cond, result) in enumerate(results):
    y = Inches(1.7 + i * 1.3)
    rect(s, L, y, W, Inches(1.1), OFF_WHITE)
    tb = box(s, Inches(1.5), y + Inches(0.1), Inches(10.3), Inches(0.9))
    txt(tb.text_frame, cond, sz=15, color=GRAY)
    para(tb.text_frame, result, sz=18, bold=True, color=BLUE, sp=Pt(4))

rect(s, L, Inches(5.7), W, Inches(1.0), RGBColor(0xE8, 0xEE, 0xF8))
tb = box(s, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.8))
txt(
    tb.text_frame,
    "Headline: up to 1.4x throughput at 20% lower cost, or 2.35x throughput at same budget",
    sz=18,
    bold=True,
    color=BLUE,
    align=PP_ALIGN.CENTER,
)

tb = box(s, L, Inches(6.8), W, Inches(0.5))
txt(
    tb.text_frame,
    "Robust to workload changes (7% drop with mismatched workload) and model changes.",
    sz=14,
    color=GRAY,
)

fig_slide(
    "Splitwise Throughput-Latency Results",
    "sw_fig16_latency.png",
    "Source: Patel et al., ISCA 2024, Figure 16",
)

# ====== SLIDE 14: SPLITWISE SUMMARY ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "Splitwise Summary", sz=34, bold=True, color=BLACK)

tb = box(s, L, Inches(1.8), W, Inches(3.5))
txt(tb.text_frame, "", sz=16, color=DARK)
for item in [
    "Characterized real Azure production workloads, found 7 key insights",
    "Designed three-pool architecture (prompt, token, mixed)",
    "Optimized KV-cache transfer: 0.8% of end-to-end latency",
    "Explored four heterogeneous cluster configurations",
    "Result: cheaper GPUs for token generation work just as well",
]:
    bullet(tb.text_frame, item, sz=17, color=DARK, sp=Pt(12))

rect(s, L, Inches(5.2), W, Inches(0.6), OFF_WHITE)
tb = box(s, Inches(1.5), Inches(5.25), Inches(10.3), Inches(0.5))
txt(
    tb.text_frame,
    "Limitations: needs cluster provisioning, less useful with few GPUs, basic fault tolerance",
    sz=14,
    color=GRAY,
    align=PP_ALIGN.CENTER,
)

# ====== SLIDE 15: DISTSERVE DIVIDER ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
rect(s, Inches(0), Inches(3.6), Inches(13.333), Pt(2), TEAL)

tb = box(s, L, Inches(2.2), W, Inches(1))
txt(tb.text_frame, "DistServe", sz=48, bold=True, color=WHITE)

tb = box(s, L, Inches(4.0), W, Inches(1.2))
txt(
    tb.text_frame,
    "Disaggregating Prefill and Decoding for Goodput-Optimized LLM Serving",
    sz=20,
    color=RGBColor(0x88, 0xCC, 0xBB),
)
para(
    tb.text_frame,
    "Zhong et al. | Peking University + StepFun + UC San Diego | OSDI 2024",
    sz=16,
    color=LIGHT_GRAY,
    sp=Pt(10),
)

# ====== SLIDE 16: GOODPUT ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(
    tb.text_frame,
    "Goodput: The Right Optimization Target",
    sz=34,
    bold=True,
    color=BLACK,
)

rect(s, L, Inches(1.7), W, Inches(1.2), RGBColor(0xE8, 0xEE, 0xF8))
tb = box(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.0))
txt(
    tb.text_frame,
    "Per-GPU Goodput = max request rate while meeting SLO targets",
    sz=20,
    bold=True,
    color=BLUE,
    align=PP_ALIGN.CENTER,
)
para(
    tb.text_frame,
    "Higher goodput = lower cost per query = what production systems actually optimize for",
    sz=16,
    color=GRAY,
    align=PP_ALIGN.CENTER,
    sp=Pt(8),
)

tb = box(s, L, Inches(3.3), W, Inches(3))
txt(tb.text_frame, "DistServe's approach:", sz=20, bold=True, color=BLACK)
bullet(
    tb.text_frame,
    "Separate prefill and decoding onto different GPU instances",
    sz=17,
    color=DARK,
    sp=Pt(14),
)
bullet(
    tb.text_frame,
    "Optimize parallelism strategy for each phase independently",
    sz=17,
    color=DARK,
)
bullet(
    tb.text_frame,
    "Automatically find the best placement on the cluster",
    sz=17,
    color=DARK,
)
bullet(
    tb.text_frame,
    "Use simulation to evaluate configurations (under 2% error)",
    sz=17,
    color=DARK,
)

# ====== SLIDE 17: TRADEOFF ANALYSIS ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(
    tb.text_frame,
    "Tradeoff Analysis After Disaggregation",
    sz=34,
    bold=True,
    color=BLACK,
)

# Prefill
rect(s, L, Inches(1.7), Inches(5.2), Inches(3.8), RGBColor(0xE8, 0xEE, 0xF8))
tb = box(s, Inches(1.5), Inches(1.8), Inches(4.6), Inches(3.6))
txt(tb.text_frame, "Prefill Instances", sz=20, bold=True, color=BLUE)
bullet(tb.text_frame, "Modeled as M/D/1 queue", sz=15, color=DARK, sp=Pt(10))
bullet(tb.text_frame, "Avg TTFT = execution time + queuing delay", sz=15, color=DARK)
bullet(tb.text_frame, "Low request rates:", sz=15, color=DARK, sp=Pt(10))
bullet(
    tb.text_frame,
    "Tensor parallelism better (cuts execution time)",
    sz=13,
    color=GRAY,
    level=1,
)
bullet(tb.text_frame, "High request rates:", sz=15, color=DARK, sp=Pt(6))
bullet(
    tb.text_frame,
    "Pipeline parallelism better (handles queuing)",
    sz=13,
    color=GRAY,
    level=1,
)

# Decode
rect(s, Inches(6.9), Inches(1.7), Inches(5.2), Inches(3.8), RGBColor(0xE8, 0xF5, 0xF0))
tb = box(s, Inches(7.2), Inches(1.8), Inches(4.6), Inches(3.6))
txt(tb.text_frame, "Decoding Instances", sz=20, bold=True, color=TEAL)
bullet(tb.text_frame, "Single job is bandwidth-bound", sz=15, color=DARK, sp=Pt(10))
bullet(tb.text_frame, "Batching is essential for utilization", sz=15, color=DARK)
bullet(tb.text_frame, "Disaggregation helps:", sz=15, color=DARK, sp=Pt(10))
bullet(
    tb.text_frame,
    "Multiple prefill instances feed one decoder",
    sz=13,
    color=GRAY,
    level=1,
)
bullet(
    tb.text_frame,
    "Naturally larger batches, better GPU use",
    sz=13,
    color=GRAY,
    level=1,
)
bullet(
    tb.text_frame, "Large batches approach compute-bound", sz=15, color=DARK, sp=Pt(6)
)

rect(s, L, Inches(5.8), W, Inches(0.7), OFF_WHITE)
tb = box(s, Inches(1.5), Inches(5.85), Inches(10.3), Inches(0.6))
txt(
    tb.text_frame,
    "After disaggregation, each phase has independent knobs. Impossible when sharing GPUs.",
    sz=17,
    bold=True,
    color=BLUE,
    align=PP_ALIGN.CENTER,
)

fig_slide(
    "DistServe System Architecture",
    "ds_fig6_architecture.png",
    "Source: Zhong et al., OSDI 2024, Figure 6",
)

# ====== SLIDE 18: PLACEMENT ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "Placement Algorithms", sz=34, bold=True, color=BLACK)

tb = box(s, L, Inches(1.4), W, Inches(0.5))
txt(
    tb.text_frame,
    "Input: model, workload, SLOs, cluster hardware.    Output: optimal parallelism + instance counts + placement.",
    sz=15,
    color=GRAY,
)

# Alg 1
rect(s, L, Inches(2.2), Inches(5.2), Inches(3.0), OFF_WHITE)
tb = box(s, Inches(1.5), Inches(2.3), Inches(4.6), Inches(2.8))
txt(tb.text_frame, "Algorithm 1: Fast cross-node network", sz=17, bold=True, color=BLUE)
bullet(tb.text_frame, "Enumerate parallelism configs", sz=14, color=DARK, sp=Pt(8))
bullet(tb.text_frame, "Simulate goodput for each", sz=14, color=DARK)
bullet(tb.text_frame, "Pick best config per phase", sz=14, color=DARK)
bullet(tb.text_frame, "Replicate to meet target rate", sz=14, color=DARK)
bullet(tb.text_frame, "Runs in under 1.3 minutes", sz=14, color=GREEN, bold=True)

# Alg 2
rect(s, Inches(6.9), Inches(2.2), Inches(5.2), Inches(3.0), OFF_WHITE)
tb = box(s, Inches(7.2), Inches(2.3), Inches(4.6), Inches(2.8))
txt(tb.text_frame, "Algorithm 2: Limited bandwidth", sz=17, bold=True, color=TEAL)
bullet(tb.text_frame, "Prefill + decode on same node", sz=14, color=DARK, sp=Pt(8))
bullet(tb.text_frame, "Use NVLINK for KV transfer", sz=14, color=DARK)
bullet(tb.text_frame, "Co-optimize within node budget", sz=14, color=DARK)
bullet(tb.text_frame, "Group layers into segments", sz=14, color=DARK)

rect(s, L, Inches(5.5), W, Inches(0.7), RGBColor(0xE8, 0xEE, 0xF8))
tb = box(s, Inches(1.5), Inches(5.55), Inches(10.3), Inches(0.6))
txt(
    tb.text_frame,
    "Simulator accuracy: under 2% error compared to real hardware runs",
    sz=16,
    bold=True,
    color=BLUE,
    align=PP_ALIGN.CENTER,
)

# ====== SLIDE 19: DISTSERVE RESULTS ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "DistServe Results", sz=34, bold=True, color=BLACK)

tb = box(s, L, Inches(1.3), W, Inches(0.5))
txt(
    tb.text_frame,
    "32 NVIDIA A100-80GB GPUs across 4 nodes  |  OPT models (13B, 66B, 175B)  |  vs vLLM and DeepSpeed-MII",
    sz=14,
    color=GRAY,
)

table = tbl(s, 4, 4, Inches(1.2), Inches(2.0), Inches(10.9), Inches(2.5))
for i, v in enumerate(
    ["Workload", "Rate vs vLLM", "SLO vs vLLM", "Rate vs DeepSpeed-MII"]
):
    cell(table, 0, i, v, sz=13, bold=True, color=WHITE, fill=RGBColor(0x55, 0x55, 0x55))
for r, row in enumerate(
    [
        [
            "Chatbot (ShareGPT)",
            "2.0x to 4.6x higher",
            "1.8x to 3.2x tighter",
            "1.6x to 7.4x higher",
        ],
        ["Code Completion", "5.7x higher", "1.4x tighter", "1.6x higher"],
        ["Summarization", "4.3x higher", "12.6x tighter", "1.8x higher"],
    ]
):
    f = OFF_WHITE if r % 2 == 0 else WHITE
    for c, v in enumerate(row):
        b = "12.6" in v or "7.4" in v or "5.7" in v
        cell(table, r + 1, c, v, sz=13, color=GREEN if b else DARK, fill=f, bold=b)

# Parallelism chosen
tb = box(s, L, Inches(4.8), W, Inches(0.5))
txt(
    tb.text_frame,
    "Parallelism strategies chosen (different per phase):",
    sz=16,
    bold=True,
    color=BLACK,
)

table2 = tbl(s, 4, 3, Inches(3.5), Inches(5.3), Inches(6.3), Inches(1.6))
for i, v in enumerate(["Model", "Prefill (TP, PP)", "Decode (TP, PP)"]):
    cell(
        table2, 0, i, v, sz=12, bold=True, color=WHITE, fill=RGBColor(0x55, 0x55, 0x55)
    )
for r, row in enumerate(
    [
        ["OPT-13B", "(2, 1)", "(1, 1)"],
        ["OPT-66B", "(4, 1)", "(2, 2)"],
        ["OPT-175B", "(3, 3)", "(4, 3)"],
    ]
):
    f = OFF_WHITE if r % 2 == 0 else WHITE
    for c, v in enumerate(row):
        cell(table2, r + 1, c, v, sz=12, color=BLUE if c > 0 else DARK, fill=f)

fig_slide(
    "DistServe Evaluation: Chatbot and Workloads",
    "ds_fig8_chatbot.png",
    "Source: Zhong et al., OSDI 2024, Figure 8",
)

fig_slide(
    "DistServe Evaluation: Code and Summarization",
    "ds_fig9_code_summ.png",
    "Source: Zhong et al., OSDI 2024, Figure 9",
)

# ====== SLIDE 20: ABLATION ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "What Actually Matters?", sz=34, bold=True, color=BLACK)

# KV overhead
rect(s, L, Inches(1.7), Inches(5.2), Inches(2.8), RGBColor(0xEF, 0xFA, 0xF4))
tb = box(s, Inches(1.5), Inches(1.8), Inches(4.6), Inches(2.6))
txt(tb.text_frame, "KV-Cache Transfer Cost", sz=20, bold=True, color=GREEN)
bullet(
    tb.text_frame, "OPT-175B: under 0.1% of total latency", sz=15, color=DARK, sp=Pt(10)
)
bullet(tb.text_frame, "95% of requests: under 30ms delay", sz=15, color=DARK)
para(
    tb.text_frame,
    "Transfer is not the bottleneck.",
    sz=16,
    bold=True,
    color=GREEN,
    sp=Pt(12),
)
para(tb.text_frame, "Confirms Splitwise's findings.", sz=14, color=GRAY, sp=Pt(4))

# Ablation
rect(s, Inches(6.9), Inches(1.7), Inches(5.2), Inches(2.8), RGBColor(0xFD, 0xF0, 0xEF))
tb = box(s, Inches(7.2), Inches(1.8), Inches(4.6), Inches(2.6))
txt(tb.text_frame, "Ablation: Parallelism Search", sz=20, bold=True, color=RED_ACCENT)
bullet(
    tb.text_frame,
    "vLLM++ (exhaustive parallelism search)",
    sz=15,
    color=DARK,
    sp=Pt(10),
)
bullet(
    tb.text_frame,
    "Same performance as default vLLM",
    sz=15,
    color=RED_ACCENT,
    bold=True,
)
para(
    tb.text_frame,
    "Interference cancels parallelism gains.",
    sz=15,
    color=DARK,
    sp=Pt(12),
)
para(tb.text_frame, "Must disaggregate first.", sz=16, bold=True, color=BLUE, sp=Pt(4))

fig_slide(
    "Per-Request Latency Distribution",
    "ds_fig10_latency.png",
    "Source: Zhong et al., OSDI 2024, Figure 10",
)

# ====== SLIDE 21: COMPARISON ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "Splitwise vs DistServe", sz=34, bold=True, color=BLACK)

table = tbl(s, 8, 3, Inches(1.2), Inches(1.5), Inches(10.9), Inches(4.5))
for i, v in enumerate(["", "Splitwise", "DistServe"]):
    cell(
        table,
        0,
        i,
        v,
        sz=14,
        bold=True,
        color=WHITE,
        fill=BLUE if i == 1 else (TEAL if i == 2 else RGBColor(0x55, 0x55, 0x55)),
    )
for r, row in enumerate(
    [
        ["Focus", "Heterogeneous hardware", "Goodput optimization"],
        ["Cluster", "Mixed GPU types", "Homogeneous"],
        ["Core idea", "Right hardware per phase", "Right parallelism per phase"],
        ["Audience", "Cloud providers", "Service operators"],
        ["KV transfer", "Per-layer overlapped", "Pull-based on-demand"],
        ["Analysis", "Production traces", "Queueing theory"],
        ["Venue", "ISCA 2024", "OSDI 2024"],
    ]
):
    f = OFF_WHITE if r % 2 == 0 else WHITE
    for c, v in enumerate(row):
        cell(table, r + 1, c, v, sz=13, color=DARK, fill=f, bold=(c == 0))

tb = box(s, Inches(1.2), Inches(6.2), Inches(10.9), Inches(0.6))
txt(
    tb.text_frame,
    "Complementary approaches. Best system would combine both.",
    sz=18,
    bold=True,
    color=BLUE,
    align=PP_ALIGN.CENTER,
)

# ====== SLIDE 22: FUTURE ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "Limitations and Open Questions", sz=34, bold=True, color=BLACK)

# Limitations
rect(s, L, Inches(1.7), Inches(5.2), Inches(2.8), RGBColor(0xFD, 0xF0, 0xEF))
tb = box(s, Inches(1.5), Inches(1.8), Inches(4.6), Inches(2.6))
txt(tb.text_frame, "Shared limitations", sz=18, bold=True, color=RED_ACCENT)
bullet(tb.text_frame, "Less beneficial with few GPUs", sz=15, color=DARK, sp=Pt(10))
bullet(tb.text_frame, "Basic fault tolerance", sz=15, color=DARK)
bullet(tb.text_frame, "FCFS can cause convoy effects", sz=15, color=DARK)
bullet(tb.text_frame, "Batch workloads may prefer colocation", sz=15, color=DARK)

# Open
rect(s, Inches(6.9), Inches(1.7), Inches(5.2), Inches(2.8), RGBColor(0xE8, 0xEE, 0xF8))
tb = box(s, Inches(7.2), Inches(1.8), Inches(4.6), Inches(2.6))
txt(tb.text_frame, "Open questions", sz=18, bold=True, color=BLUE)
bullet(tb.text_frame, "Preemptive scheduling", sz=15, color=DARK, sp=Pt(10))
bullet(tb.text_frame, "Combining heterogeneous HW + placement", sz=15, color=DARK)
bullet(tb.text_frame, "Million-token contexts", sz=15, color=DARK)
bullet(tb.text_frame, "Interaction with MoE models", sz=15, color=DARK)

rect(s, L, Inches(4.9), W, Inches(1.0), OFF_WHITE)
tb = box(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.8))
txt(
    tb.text_frame,
    "Already adopted by SGLang, vLLM, and Mooncake.",
    sz=17,
    bold=True,
    color=BLACK,
    align=PP_ALIGN.CENTER,
)
para(
    tb.text_frame,
    "Phase disaggregation is becoming the standard for production LLM serving.",
    sz=16,
    color=GRAY,
    align=PP_ALIGN.CENTER,
    sp=Pt(4),
)

# ====== SLIDE 23: SUMMARY ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "Summary", sz=36, bold=True, color=BLACK)

# Splitwise
rect(s, L, Inches(1.7), Inches(5.2), Inches(2.5), RGBColor(0xE8, 0xEE, 0xF8))
tb = box(s, Inches(1.5), Inches(1.8), Inches(4.6), Inches(2.3))
txt(tb.text_frame, "Splitwise (ISCA 2024)", sz=20, bold=True, color=BLUE)
bullet(
    tb.text_frame, "Production workload characterization", sz=15, color=DARK, sp=Pt(8)
)
bullet(tb.text_frame, "Heterogeneous hardware per phase", sz=15, color=DARK)
bullet(
    tb.text_frame, "1.4x throughput at 20% lower cost", sz=15, color=GREEN, bold=True
)
bullet(tb.text_frame, "or 2.35x throughput, same budget", sz=15, color=GREEN, bold=True)

# DistServe
rect(s, Inches(6.9), Inches(1.7), Inches(5.2), Inches(2.5), RGBColor(0xE8, 0xF5, 0xF0))
tb = box(s, Inches(7.2), Inches(1.8), Inches(4.6), Inches(2.3))
txt(tb.text_frame, "DistServe (OSDI 2024)", sz=20, bold=True, color=TEAL)
bullet(tb.text_frame, "Goodput-optimal placement", sz=15, color=DARK, sp=Pt(8))
bullet(tb.text_frame, "Automatic parallelism optimization", sz=15, color=DARK)
bullet(tb.text_frame, "Up to 7.4x higher request rate", sz=15, color=GREEN, bold=True)
bullet(tb.text_frame, "Up to 12.6x tighter SLO", sz=15, color=GREEN, bold=True)

rect(s, L, Inches(4.6), W, Inches(1.2), OFF_WHITE)
tb = box(s, Inches(1.5), Inches(4.7), Inches(10.3), Inches(1.0))
txt(
    tb.text_frame,
    "Prefill and decoding are different workloads.",
    sz=20,
    color=BLACK,
    align=PP_ALIGN.CENTER,
)
para(
    tb.text_frame,
    "Treating them as one wastes hardware, power, and money.",
    sz=20,
    bold=True,
    color=BLUE,
    align=PP_ALIGN.CENTER,
    sp=Pt(4),
)

# ====== REFERENCES SLIDE ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

tb = box(s, L, T_TITLE, W, Inches(0.7))
txt(tb.text_frame, "References", sz=34, bold=True, color=BLACK)

refs = [
    "[1]  P. Patel, E. Choukse, C. Zhang, A. Shah, I. Goiri, S. Maleki, R. Bianchini.",
    '     "Splitwise: Efficient Generative LLM Inference Using Phase Splitting."',
    "     ISCA 2024. arXiv:2311.18677",
    "",
    "[2]  Y. Zhong, S. Liu, J. Chen, J. Hu, Y. Zhu, X. Liu, X. Jin, H. Zhang.",
    '     "DistServe: Disaggregating Prefill and Decoding for Goodput-optimized',
    '     Large Language Model Serving."  OSDI 2024. arXiv:2401.09670',
    "",
    "[3]  A. Agrawal et al. Sarathi: Efficient LLM Inference by Piggybacking Decodes",
    "     with Chunked Prefills. arXiv:2308.16369, 2023.",
    "",
    "[4]  W. Kwon et al. Efficient Memory Management for Large Language Model Serving",
    "     with PagedAttention (vLLM). SOSP 2023.",
    "",
    "[5]  Y. Sheng et al. FlexGen: High-Throughput Generative Inference of Large Language",
    "     Models with a Single GPU. ICML 2023.",
]

tb = box(s, L, Inches(1.5), Inches(10.5), Inches(5.5))
txt(tb.text_frame, refs[0], sz=13, color=DARK)
for line in refs[1:]:
    para(tb.text_frame, line, sz=13, color=DARK if line else DARK, sp=Pt(1))

# ====== SLIDE 24: THANK YOU ======
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
rect(s, Inches(0), Inches(3.2), Inches(13.333), Pt(2), BLUE)

tb = box(s, L, Inches(1.8), W, Inches(1))
txt(tb.text_frame, "Thank You", sz=48, bold=True, color=WHITE)

tb = box(s, L, Inches(3.6), W, Inches(0.7))
txt(tb.text_frame, "Questions?", sz=32, color=RGBColor(0x88, 0xAA, 0xDD))

tb = box(s, L, Inches(4.8), W, Inches(2))
txt(tb.text_frame, "Splitwise: arXiv:2311.18677 (ISCA 2024)", sz=16, color=LIGHT_GRAY)
para(
    tb.text_frame,
    "DistServe: arXiv:2401.09670 (OSDI 2024)",
    sz=16,
    color=LIGHT_GRAY,
    sp=Pt(4),
)
para(tb.text_frame, "", sz=10, color=LIGHT_GRAY, sp=Pt(16))
para(
    tb.text_frame,
    "Berat Celik  &  Jiayang (Ethan) Chen",
    sz=20,
    bold=True,
    color=WHITE,
    sp=Pt(4),
)
para(tb.text_frame, "ECE 5545  |  Spring 2026", sz=16, color=LIGHT_GRAY, sp=Pt(4))


# ====== ADD SLIDE NUMBERS ======
DARK_SLIDES = {1, 8, 19, 33}
for i, slide in enumerate(prs.slides):
    num = i + 1
    if num == 1:
        continue
    clr = RGBColor(0x55, 0x66, 0x88) if num in DARK_SLIDES else LIGHT_GRAY
    tb = box(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.35))
    txt(tb.text_frame, str(num), sz=10, color=clr, align=PP_ALIGN.RIGHT)

prs.save(
    "/Users/beratcelik/Desktop/hml presentation/presentation/Phase_Disaggregation_LLM_Inference.pptx"
)
print("DONE: 33 slides with paper figures, references, Palatino Linotype font.")
